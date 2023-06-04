#%% Generating an automatic presentation from a notebook using python-pptx, nbformat and GPT-4

import nbformat
from pptx import Presentation
from pptx.util import Inches

def read_ipynb_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        notebook = nbformat.read(file, as_version=4)
    return notebook

def create_presentation(notebook, template_path):
    prs = Presentation(template_path)
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[4]

    # Set the presentation title
    for cell in notebook.cells:
        if cell.cell_type == "markdown" and cell.source.startswith("# "):
            title_slide = prs.slides.add_slide(title_slide_layout)
            title = title_slide.shapes.title
            title.text = cell.source[2:].strip()
            break

    # Create slides for each second-level header
    for cell in notebook.cells:
        if cell.cell_type == "markdown" and cell.source.startswith("## "):
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            lines = cell.source.split('\n')
            title.text = lines[0][3:].strip()

            # Add content to the slide
            content = ""
            for line in lines[1:]:
                content += line.strip() + "\n"

            for subcell in notebook.cells[notebook.cells.index(cell) + 1:]:
                if subcell.cell_type == "markdown" and subcell.source.startswith("### "):
                    content += subcell.source[4:].strip() + "\n\n"
                elif subcell.cell_type == "markdown" and subcell.source.startswith("## "):
                    break
                elif subcell.cell_type == "markdown":
                    lines = subcell.source.split('\n')
                    for line in lines:
                        if not line.startswith("### "):
                            content += line.strip() + "\n"

            # Find the existing text shape on the slide and insert the content
            for shape in slide.shapes:
                if shape.has_text_frame and not shape.text.startswith("Click to edit"):
                    text_frame = shape.text_frame
                    text_frame.text = content
                    break

    return prs

#%% ------------------------------- PATHS ------------------------------- #

notebook_path = "/Users/nilsjennissen/PycharmProjects/presentations/notebooks/template.ipynb"
template_path = "../templates/template.pptx"
output_path = "../pres/output.pptx"


notebook = read_ipynb_file(notebook_path)
presentation = create_presentation(notebook, template_path)
presentation.save(output_path)