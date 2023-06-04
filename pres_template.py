#%% Packages
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

#%% Presentation
# Define the path to the powerpoint presentation
path = 'templates/template.pptx'
# Open the powerpoint presentation
prs = Presentation(path)

#%%
# Get summary details of the presentation
number_of_slides = len(prs.slides)
print('Number of slides: ', number_of_slides)

#%%
# Create lists for slide information
slide_titles = prs.slides[0].shapes.title.text
#%%
slide_body = prs.slides[3].shapes.placeholders[1].text
slide_body_text = prs.slides[3].shapes.placeholders[1].text_frame.paragraphs[0].runs[0].text

for shape in prs.slides[3].placeholders:
     print('%d %s' % (shape.placeholder_format.idx, shape.name))


#%%
# Create a dataframe with the slide information
df = pd.DataFrame({'Slide': slide_titles, 'Body': slide_body, 'Body_text': slide_body_text})
df.head

#%%
# Return length of slide layout
print(len(prs.slide_layouts))


#%% - TITLE SLIDE -
# Create a title slide

'''Slide Introduction'''
# Create the slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add Title
title.text = "Header Text"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(73, 94, 121)

# Add Subtitle
subtitle.text = 'Subtitle text'
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(73, 94, 121)

# Add Image in the preset position
left = Inches(0.5)
top = Inches(1.5)
height = Inches(1.5)
pic = slide.shapes.add_picture('images/openAi.png', left, top, height=height)

# Add Textbox
left = Inches(0.5)
top = Inches(3.5)
width = Inches(5)
height = Inches(1)
textbox = slide.shapes.add_textbox(left, top, width, height)


#%% Save the presentation
prs.save('pres/presentation1.pptx')