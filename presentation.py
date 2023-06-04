
#%% Packages
# Import libraries
import pandas as pd
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
#from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
#from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

# import placeholders
#from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE


#%% Presentation
# Define the path to the powerpoint presentation
path = 'templates/template.pptx'
prs = Presentation(path)

for slide in prs.slides:
    print(type(slide), slide.name, 'slide_index [',prs.slides.index(slide), ']' )
    for p in slide.placeholders:
        print('index [', p.placeholder_format.idx, ']', p.placeholder_format.type, p.name)


# Get summary details of the presentation
number_of_slides = len(prs.slides)
print('Number of slides: ', number_of_slides)

#%%
# Create lists for slide headers
slide_titles = []
slide_body = []

# Extract slide text
slide_text = []
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                slide_text.append(run.text)

slide_paragraphs = []
for slide in prs.slides:
    if slide.shapes.has_text_frame:
        slide_paragraphs.append(slide.shapes[3].text)
    else:
        slide_paragraphs.append('')

#%% Inspect first slide
# First slide header
print(prs.slides[0].shapes[0].name)
print(prs.slides[0].shapes[1].name)
print(prs.slides[0].shapes[2].name)
print(prs.slides[0].shapes[3].name)
print(prs.slides[0].shapes.title.text)
print(prs.slides[0].placeholders[1])



#%%
# First shape
print(prs.slides[0].shapes.title.text)
print(prs.slides[1].shapes[0].name)

for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text_frame = shape.text_frame

#%% Map the power point object
# Create a dataframe with the shape information of each slide in a column
prs = Presentation()
for slide in prs.slides:
    print(slide.title)
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))


#%%

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[3])
for shape in slide.shapes:
    print('%s' % shape.shape_type)

#%%
slide = prs.slides.add_slide(prs.slides[3].shapes)
for shape in slide.shapes:
    print(shape.name, shape.shape_type, shape.has_text_frame, shape.is_placeholder)